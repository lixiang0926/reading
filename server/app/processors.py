from typing import Optional
import docx
from PyPDF2 import PdfReader
import io
import os
import nltk
from nltk.tokenize import sent_tokenize, word_tokenize
import asyncio
from .utils.cache import Cache
from .config import settings
from .utils.document_manager import DocumentManager
import markdown2
from bs4 import BeautifulSoup
import openpyxl
from pptx import Presentation
import ebooklib
from ebooklib import epub
import chardet
from pygments import highlight
from pygments.lexers import get_lexer_for_filename
from pygments.formatters import HtmlFormatter
import xml.etree.ElementTree as ET
import json
import csv

class FileProcessor:
    def __init__(self):
        self.cache = Cache()
        self.doc_manager = DocumentManager()
        self.processors = {
            # 文本文件
            '.txt': self.process_txt,
            '.md': self.process_markdown,
            '.csv': self.process_csv,
            '.json': self.process_json,
            '.xml': self.process_xml,
            '.html': self.process_html,
            '.htm': self.process_html,
            # Microsoft Office
            '.doc': self.process_word,
            '.docx': self.process_word,
            '.xls': self.process_excel,
            '.xlsx': self.process_excel,
            '.ppt': self.process_powerpoint,
            '.pptx': self.process_powerpoint,
            '.rtf': self.process_rtf,
            # OpenDocument
            '.odt': self.process_odt,
            '.ods': self.process_ods,
            '.odp': self.process_odp,
            # PDF
            '.pdf': self.process_pdf,
            # 电子书
            '.epub': self.process_epub,
            '.mobi': self.process_mobi,
            # 代码文件
            '.py': self.process_code,
            '.js': self.process_code,
            '.java': self.process_code,
            '.cpp': self.process_code,
            '.c': self.process_code,
            '.h': self.process_code,
            '.cs': self.process_code,
            '.php': self.process_code
        }

    async def process_file(self, file_path: str, file_ext: str, bionic_enabled: bool,
                         page: int = 1, user_id: Optional[str] = None) -> dict:
        try:
            # 获取文档ID
            doc_id = self.doc_manager.get_document_id(file_path)
            
            # 检查是否有缓存的页面
            pages_data = await self.doc_manager.get_pages(
                doc_id, 
                page - 1, 
                settings.MAX_PAGES_PER_REQUEST
            )
            
            if not pages_data:
                # 获取文件编码
                with open(file_path, 'rb') as file:
                    raw_data = file.read()
                    encoding = chardet.detect(raw_data)['encoding'] or 'utf-8'

                # 获取对应的处理器
                processor = self.processors.get(file_ext.lower())
                if not processor:
                    raise ValueError(f"不支持的文件格式: {file_ext}")

                # 处理文件
                content = await processor(file_path, encoding)

                # 分页处理
                pages = await self.doc_manager.split_content(content)
                
                # 保存分页结果
                await self.doc_manager.save_pages(doc_id, pages)
                
                # 获取请求的页面
                pages_data = await self.doc_manager.get_pages(
                    doc_id, 
                    page - 1, 
                    settings.MAX_PAGES_PER_REQUEST
                )

            # 对请求的页面应用仿生阅读处理
            if bionic_enabled:
                pages_data['pages'] = [
                    await self.apply_bionic_reading(page_content)
                    for page_content in pages_data['pages']
                ]

            # 保存阅读进度
            if user_id:
                await self.doc_manager.save_progress(doc_id, user_id, page)

            return {
                'success': True,
                'content': pages_data['pages'],
                'current_page': pages_data['current_page'],
                'total_pages': pages_data['total_pages'],
                'has_more': pages_data['has_more']
            }

        except Exception as e:
            raise Exception(f"处理文件失败: {str(e)}")

    async def process_markdown(self, file_path: str, encoding: str) -> str:
        with open(file_path, 'r', encoding=encoding) as file:
            content = file.read()
            html = markdown2.markdown(content)
            soup = BeautifulSoup(html, 'html.parser')
            return '\n'.join([f"<p>{p.get_text()}</p>" for p in soup.find_all(['p', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6'])])

    async def process_excel(self, file_path: str, encoding: str) -> str:
        def _process():
            wb = openpyxl.load_workbook(file_path)
            content = []
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                content.append(f"<h2>{sheet}</h2>")
                for row in ws.iter_rows(values_only=True):
                    content.append(f"<p>{' | '.join(str(cell) for cell in row if cell is not None)}</p>")
            return '\n'.join(content)
        return await asyncio.to_thread(_process)

    async def process_powerpoint(self, file_path: str, encoding: str) -> str:
        def _process():
            prs = Presentation(file_path)
            content = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        content.append(f"<p>{shape.text}</p>")
            return '\n'.join(content)
        return await asyncio.to_thread(_process)

    async def process_html(self, file_path: str, encoding: str) -> str:
        with open(file_path, 'r', encoding=encoding) as file:
            soup = BeautifulSoup(file.read(), 'html.parser')
            return '\n'.join([f"<p>{p.get_text()}</p>" for p in soup.find_all(['p', 'div', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6'])])

    async def process_xml(self, file_path: str, encoding: str) -> str:
        tree = ET.parse(file_path)
        root = tree.getroot()
        content = []
        for elem in root.iter():
            if elem.text and elem.text.strip():
                content.append(f"<p>{elem.text}</p>")
        return '\n'.join(content)

    async def process_json(self, file_path: str, encoding: str) -> str:
        with open(file_path, 'r', encoding=encoding) as file:
            data = json.load(file)
            return self._format_json_content(data)

    def _format_json_content(self, data, level=0):
        content = []
        if isinstance(data, dict):
            for key, value in data.items():
                if isinstance(value, (dict, list)):
                    content.append(f"<p>{'  ' * level}{key}:</p>")
                    content.append(self._format_json_content(value, level + 1))
                else:
                    content.append(f"<p>{'  ' * level}{key}: {value}</p>")
        elif isinstance(data, list):
            for item in data:
                if isinstance(item, (dict, list)):
                    content.append(self._format_json_content(item, level + 1))
                else:
                    content.append(f"<p>{'  ' * level}- {item}</p>")
        return '\n'.join(content)

    async def process_csv(self, file_path: str, encoding: str) -> str:
        content = []
        with open(file_path, 'r', encoding=encoding) as file:
            reader = csv.reader(file)
            for row in reader:
                content.append(f"<p>{' | '.join(row)}</p>")
        return '\n'.join(content)

    async def process_epub(self, file_path: str, encoding: str) -> str:
        def _process():
            book = epub.read_epub(file_path)
            content = []
            for item in book.get_items():
                if item.get_type() == ebooklib.ITEM_DOCUMENT:
                    soup = BeautifulSoup(item.get_content(), 'html.parser')
                    for p in soup.find_all(['p', 'div', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6']):
                        if p.get_text().strip():
                            content.append(f"<p>{p.get_text()}</p>")
            return '\n'.join(content)
        return await asyncio.to_thread(_process)

    async def process_code(self, file_path: str, encoding: str) -> str:
        with open(file_path, 'r', encoding=encoding) as file:
            content = file.read()
            lexer = get_lexer_for_filename(file_path)
            formatter = HtmlFormatter(style='monokai', noclasses=True)
            highlighted = highlight(content, lexer, formatter)
            return f"<div class='code'>{highlighted}</div>"

    async def process_word(self, file_path: str, encoding: str) -> str:
        def _process():
            doc = docx.Document(file_path)
            paragraphs = []
            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append(f"<p>{para.text}</p>")
            return "\n".join(paragraphs)

        return await asyncio.to_thread(_process)

    async def process_pdf(self, file_path: str, encoding: str) -> str:
        def _process():
            with open(file_path, 'rb') as file:
                reader = PdfReader(file)
                paragraphs = []
                for page in reader.pages[:settings.MAX_PDF_PAGES]:
                    text = page.extract_text()
                    if text.strip():
                        # 按行分割并过滤空行
                        lines = [line.strip() for line in text.split('\n') if line.strip()]
                        for line in lines:
                            paragraphs.append(f"<p>{line}</p>")
                return "\n".join(paragraphs)

        return await asyncio.to_thread(_process)

    async def process_txt(self, file_path: str, encoding: str) -> str:
        async def _process():
            with open(file_path, 'r', encoding=encoding) as file:
                content = file.read()
                paragraphs = []
                for line in content.split('\n'):
                    if line.strip():
                        paragraphs.append(f"<p>{line}</p>")
                return "\n".join(paragraphs)

        return await _process()

    async def apply_bionic_reading(self, content: str) -> str:
        def _process_word(word: str) -> str:
            # 跳过空白内容
            if not word.strip():
                return word

            # 处理连字符
            if "-" in word:
                return "-".join([_process_word(part) for part in word.split("-")])

            # 跳过数字
            if any(char.isdigit() for char in word):
                return word

            # 跳过特殊标点
            if all(not char.isalnum() for char in word):
                return word

            # 智能计算加粗长度
            length = len(word)
            if length <= 1:
                return word
            elif length <= 3:
                mid_point = 1
            elif length <= 6:
                mid_point = length // 2
            else:
                mid_point = (length * 3) // 5  # 60%加粗比例

            return f"<b>{word[:mid_point]}</b>{word[mid_point:]}"

        def _process():
            # 移除HTML标签
            text = content.replace('<p>', '').replace('</p>', '\n')
            
            # 处理每个句子
            sentences = sent_tokenize(text)
            processed_sentences = []
            
            for sentence in sentences:
                # 分词，保留空格和标点
                words = []
                current_word = ""
                for char in sentence:
                    if char.isspace() or char in ".,!?;:":
                        if current_word:
                            words.append(current_word)
                            current_word = ""
                        words.append(char)
                    else:
                        current_word += char
                if current_word:
                    words.append(current_word)

                # 处理每个单词
                processed_words = [_process_word(word) for word in words]
                processed_sentence = ''.join(processed_words)
                
                # 添加段落标签
                processed_sentences.append(f"<p>{processed_sentence}</p>")
            
            return '\n'.join(processed_sentences)

        return await asyncio.to_thread(_process)

# 下载必要的NLTK数据
nltk.download('punkt')
